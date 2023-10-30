import docx
from simplify_docx import simplify
from pptx import Presentation
from datasets import Dataset, load_dataset, concatenate_datasets
from os.path import splitext
import os
import pandas as pd
import PyPDF2
import zipfile
import boto3


def dataset_from_items(formatted_list):
    filtered_list = [item for item in formatted_list if item['text'] != '' or len(item["text"]) > 3]
    data = {key: [item[key] for item in filtered_list] for key in filtered_list[0]}
    dataset = Dataset.from_dict(data)
    return dataset


def dataset_from_docx(mydoc):
    # read in doc file
    my_doc = docx.Document(mydoc)

    # coerce to JSON using the standard options
    my_doc_as_json = simplify(my_doc, {"remove-leading-white-space": False})

    # extract body of document
    json_list = my_doc_as_json['VALUE'][0]['VALUE']

    # format to json dataset
    formatted_list = [
        {'label': mydoc, 'index': idx, 'text': item['VALUE'][0]['VALUE'] if item['TYPE'] == 'paragraph' else ""} for
        idx, item in enumerate(json_list)]

    return dataset_from_items(formatted_list)


def dataset_from_pptx(pptx_file):
    # Load the PowerPoint presentation
    presentation = Presentation(pptx_file)

    # Initialize an empty string to store the extracted text
    formatted_list = []

    # Iterate through slides and extract text
    idx = -1
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                idx += 1
                formatted_list.append({"label": pptx_file, "index": idx, "text": shape.text})

    return dataset_from_items(formatted_list)


def dataset_from_xlsx(file):
    # Read the Excel file into a dictionary of DataFrames, where keys are sheet names
    xls = pd.read_excel(file, sheet_name=None)
    # df = pd.read_excel(excel_file)

    formatted_list = []

    idx = -1
    for sheet_name, df in xls.items():
        # sheet_text = f"Sheet Name: {sheet_name}\n"  # Include the sheet name
        for column in df.columns:
            for value in df[column]:
                if isinstance(value, str):
                    idx += 1
                    formatted_list.append({"label": file + " " + sheet_name, "index": idx, "text": value + " "})

    return dataset_from_items(formatted_list)


def dataset_from_pdf(file):
    formatted_list = []
    with open(file, 'rb') as pdf_file:
        # Create a PDF reader object
        pdf_reader = PyPDF2.PdfReader(pdf_file)

        # Iterate through each page of the PDF
        idx = -1
        for page_num in range(len(pdf_reader.pages)):
            # Extract the text from the page
            page = pdf_reader.pages[page_num]
            idx += 1
            formatted_list.append({"label": file, "index": idx, "text": page.extract_text()})

    return dataset_from_items(formatted_list)


# private mapping
__mapping = {
    ".doc": dataset_from_docx,  # this actually does not work. doc file can not be loaded directly
    ".docx": dataset_from_docx,
    ".ppt": dataset_from_pptx,
    ".pptx": dataset_from_pptx,
    ".xls": dataset_from_xlsx,
    ".xlsx": dataset_from_xlsx,
    ".pdf": dataset_from_pdf,
}


def dataset_from_file(path):
    root, ext = splitext(path)
    return __mapping[ext](path)


def dataset_from_folder(directory_path):
    result_list = []
    for root, dirs, files in os.walk(directory_path):
        for file in files:
            file_path = os.path.join(root, file)
            print("processing:" + file_path)
            result = dataset_from_file(file_path)
            result_list.append(result)

    for dir_name in dirs:
        subdir_path = os.path.join(root, dir_name)
        result_list += dataset_from_folder(subdir_path)

    merged_dataset = concatenate_datasets(result_list)
    return merged_dataset


def dataset_preprocess(source_folder, dataset_save_path):
    dataset = dataset_from_folder(source_folder)
    dataset.save_to_disk(dataset_save_path)


def zip_folder(src_folder, output_filename):
    """Zip a folder of files
    Line 2 of comment...
    And so on... 
    """
    with zipfile.ZipFile(output_filename, 'w', zipfile.ZIP_STORED) as zipf:
        for root, dirs, files in os.walk(src_folder):
            for file in files:
                file_path = os.path.join(root, file)
                arcname = os.path.relpath(file_path, src_folder)
                zipf.write(file_path, arcname)


def unzip_folder(zip_file, folder_path):
    # Extract the zip file
    with zipfile.ZipFile(zip_file, 'r') as zip_ref:
        zip_ref.extractall(folder_path)


def model_to_s3(model_path, s3_bucket_name):
    basename = os.path.basename(model_path)

    zip_filename = basename + ".zip"
    zip_folder(model_path, zip_filename)

    # Upload the zipped file to S3
    # s3 = boto3.client('s3', aws_access_key_id=aws_access_key_id, aws_secret_access_key=aws_secret_access_key)
    s3 = boto3.client('s3')
    s3.upload_file(zip_filename, s3_bucket_name, zip_filename)
    # TODO: for later
    # os.remove(zip_filename)


def model_from_s3(s3_path):
    s3 = boto3.client('s3')

    s3_bucket = os.path.dirname(s3_path)
    zip_file = os.path.basename(s3_path)

    # Download the zip file from S3
    s3.download_file(s3_bucket, zip_file, zip_file)
    local_path = os.path.splitext(zip_file)[0]
    unzip_folder(zip_file, local_path)

    # Clean up the downloaded zip file
    os.remove(zip_file)
