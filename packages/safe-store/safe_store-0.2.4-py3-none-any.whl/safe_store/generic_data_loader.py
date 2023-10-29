from pathlib import Path
from typing import List
 
class GenericDataLoader:
    @staticmethod
    def read_file(file_path: Path) -> str:
        """
        Read a file and return its content as a string.

        Args:
            file_path (Path): The path to the file.

        Returns:
            str: The content of the file.
        
        Raises:
            ValueError: If the file type is unknown.
        """
        if file_path.suffix ==".pdf":
            return GenericDataLoader.read_pdf_file(file_path)
        elif file_path.suffix == ".docx":
            return GenericDataLoader.read_docx_file(file_path)
        elif file_path.suffix == ".json":
            return GenericDataLoader.read_json_file(file_path)
        elif file_path.suffix == ".html":
            return GenericDataLoader.read_html_file(file_path)
        elif file_path.suffix == ".pptx":
            return GenericDataLoader.read_pptx_file(file_path)
        if file_path.suffix in [".txt", ".rtf", ".md", ".log", ".cpp", ".java", ".js", ".py", ".rb", ".sh", ".sql", ".css", ".html", ".php", ".json", ".xml", ".yaml", ".yml", ".h", ".hh", ".hpp", ".inc", ".snippet", ".snippets", ".asm", ".s", ".se", ".sym", ".ini", ".inf", ".map", ".bat"]:
            return GenericDataLoader.read_text_file(file_path)
        else:
            raise ValueError("Unknown file type")
        
    @staticmethod
    def get_supported_file_types() -> List[str]:
        """
        Get the list of supported file types.

        Returns:
            List[str]: The list of supported file types.
        """
        return ["pdf", "txt", "docx", "json", "html", "pptx",".txt", ".md", ".log", ".cpp", ".java", ".js", ".py", ".rb", ".sh", ".sql", ".css", ".html", ".php", ".json", ".xml", ".yaml", ".yml", ".h", ".hh", ".hpp", ".inc", ".snippet", ".snippets", ".asm", ".s", ".se", ".sym", ".ini", ".inf", ".map", ".bat", ".rtf"]    
    
    @staticmethod
    def read_pdf_file(file_path: Path) -> str:
        """
        Read a PDF file and return its content as a string.

        Args:
            file_path (Path): The path to the PDF file.

        Returns:
            str: The content of the PDF file.
        """
        import PyPDF2
        def extract_text_from_pdf(file_path):
            text = ""
            with open(file_path, 'rb') as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                for page in pdf_reader.pages:
                    text += page.extract_text()
            return text
        # Extract text from the PDF
        text = extract_text_from_pdf(file_path)

        # Convert to Markdown (You may need to implement custom logic based on your specific use case)
        markdown_text = text.replace('\n', '  \n')  # Adding double spaces at the end of each line for Markdown line breaks
        
        return markdown_text

    @staticmethod
    def read_docx_file(file_path: Path) -> str:
        """
        Read a DOCX file and return its content as a string.

        Args:
            file_path (Path): The path to the DOCX file.

        Returns:
            str: The content of the DOCX file.
        """
        try:
            from docx import Document
        except ImportError:
            PackageManager.install_package("python-docx")
            from docx import Document
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    @staticmethod
    def read_json_file(file_path: Path) -> str:
        """
        Read a JSON file and return its content as a string.

        Args:
            file_path (Path): The path to the JSON file.

        Returns:
            str: The content of the JSON file.
        """
        import json
        with open(file_path, 'r', encoding='utf-8') as file:
            data = str(json.load(file))
        return data
    
    @staticmethod
    def read_csv_file(file_path: Path) -> str:
        """
        Read a CSV file and return its content as a string.
        Args:
            file_path (Path): The path to the CSV file.
        Returns:
            str: The content of the CSV file.
        """
        with open(file_path, 'r') as file:
            content = file.read()
        return content   

    @staticmethod
    def read_html_file(file_path: Path) -> str:
        """
        Read an HTML file and return its content as a string.

        Args:
            file_path (Path): The path to the HTML file.

        Returns:
            str: The content of the HTML file.
        """
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            PackageManager.install_package("beautifulsoup4")
            from bs4 import BeautifulSoup
        with open(file_path, 'r') as file:
            soup = BeautifulSoup(file, 'html.parser')
            text = soup.get_text()
        return text
    
    @staticmethod
    def read_pptx_file(file_path: Path) -> str:
        """
        Read a PPTX file and return its content as a string.

        Args:
            file_path (Path): The path to the PPTX file.

        Returns:
            str: The content of the PPTX file.
        """
        try:
            from pptx import Presentation
        except ImportError:
            PackageManager.install_package("python-pptx")
            from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        return text
    
    @staticmethod
    def read_text_file(file_path: Path) -> str:
        """
        Read a text file and return its content as a string.

        Args:
            file_path (Path): The path to the text file.

        Returns:
            str: The content of the text file.
        """
        # Implementation details omitted for brevity
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        return content

